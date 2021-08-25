import openpyxl
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

def start():
    # excel_file = input('What is the path of the excel file?')
    wb_obj = openpyxl.load_workbook('names.xlsx')
    sheet_obj = wb_obj.active
    for i in range(7, 94):
        cell_obj_1 = sheet_obj.cell(i, 1)
        cell_obj_2 = sheet_obj.cell(i, 2)
        name = cell_obj_1.value + cell_obj_2.value
        print(name)
        powerpoint(name)

def powerpoint(name):
    prs = Presentation('Raúl Antonio Aguilar Vera.pptx')
    slide = prs.slides[0]
    paragraphs = slide.shapes[3].text_frame.paragraphs
    for paragraph in paragraphs:
        for run in paragraph.runs:
            run.text = name
            print('name inserted')
            paragraph.alignment = PP_ALIGN.CENTER
            print(paragraph.alignment)
    print('filename = ' + name + '.pptx')
    prs.save(name + '.pptx')

start()